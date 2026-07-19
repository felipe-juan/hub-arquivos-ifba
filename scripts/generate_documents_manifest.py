#!/usr/bin/env python3
"""
Generate documents/manifest.json for HUB Arquivos IFBA.

This version is intentionally safe for repeated use:
- It rewrites the manifest from scratch every run, so old entries cannot accumulate.
- It ignores helper/mapping/readme files that should never appear in the public hub.
- It removes duplicates from the generated manifest by file hash, so duplicated files do not
  show twice on the website.
- It writes a duplicate report, but it does NOT delete your files.

Usage from the project root:
  python3 scripts/generate_documents_manifest.py

Recommended optional dependencies:
  python3 -m pip install --user pymupdf pypdf openpyxl python-docx python-pptx

Supported searchable extraction:
  PDF: PyMuPDF, pypdf, pdftotext fallback
  TXT/MD: native
  DOCX/XLSX/PPTX: optional Python packages

For scanned image PDFs, add OCR first or create a same-name .txt sidecar.
Example:
  documents/ppcs/ppc-2024.pdf
  documents/ppcs/ppc-2024.txt
"""
from __future__ import annotations

import csv
import hashlib
import json
import re
import shutil
import subprocess
import zipfile
import xml.etree.ElementTree as ET
import sys
import unicodedata
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "documents"
MANIFEST_JSON = DOCS_DIR / "manifest.json"
MANIFEST_SUMMARY_JSON = DOCS_DIR / "manifest-summary.json"
SEARCH_INDEX_JSON = DOCS_DIR / "search-index.json"
MANIFEST_CSV = DOCS_DIR / "manifest.csv"
REPORT_DIR = DOCS_DIR / "_manifest_reports"
THUMBNAIL_DIR = ROOT / "assets" / "document-thumbnails"
THUMBNAIL_SIZES = (160, 320, 520)
DUPLICATES_CSV = REPORT_DIR / "duplicates-ignored.csv"
IGNORED_CSV = REPORT_DIR / "ignored-files.csv"

SUPPORTED = {
    ".pdf": "PDF",
    ".png": "Imagem",
    ".jpg": "Imagem",
    ".jpeg": "Imagem",
    ".webp": "Imagem",
    ".gif": "Imagem",
    ".txt": "Texto",
    ".md": "Texto",
    ".doc": "Word/Writer",
    ".docx": "Word/Writer",
    ".odt": "Word/Writer",
    ".xls": "Excel/Calc",
    ".xlsx": "Excel/Calc",
    ".ods": "Excel/Calc",
    ".ppt": "PowerPoint/Impress",
    ".pptx": "PowerPoint/Impress",
    ".odp": "PowerPoint/Impress",
}

PRIMARY_EXTENSIONS = set(SUPPORTED) - {".txt", ".md"}

CATEGORY_LABELS = {
    "ppcs": "PPCs",
    "ppc": "PPCs",
    "matrizes-curriculares": "Matrizes curriculares",
    "matrizes": "Matrizes curriculares",
    "regulamentos-bsi": "Regulamentos BSI",
    "regulamentos": "Regulamentos",
    "portarias": "Portarias",
    "normas-ifba": "Normas IFBA",
    "normas": "Normas IFBA",
    "diretrizes-cne": "Diretrizes CNE/MEC",
    "diretrizes": "Diretrizes CNE/MEC",
    "resolucoes": "Resoluções",
    "resoluções": "Resoluções",
    "infraestrutura": "Infraestrutura",
    "coordenacao": "Coordenação de Curso",
    "coordenação": "Coordenação de Curso",
    "colegiado": "Colegiado de Curso",
    "nde": "Núcleo Docente Estruturante",
}

KNOWN_CATEGORY_FOLDERS = set(CATEGORY_LABELS)

ACRONYMS = {
    "ppc", "nde", "cne", "ces", "consepe", "consup", "concam", "ifba",
    "bsi", "si", "vca", "ldb", "sinaes", "pne", "tea", "libras", "tcc",
    "naes", "conaes", "mec", "ace", "ead", "pdf", "dg"
}

# Files matching these patterns must never appear as public documents.
# Matching is accent-insensitive and case-insensitive.
IGNORE_NAME_PATTERNS = [
    r"^leia[\s_-]*me.*",
    r".*renomeacao.*",
    r".*renomeacao.*",  # accent-normalized duplicate kept for clarity
    r".*mapping.*",
    r"^readme.*",
    r"^manifest\.csv$",
    r"^manifest\.json$",
    r"^content-blueprint\.json$",
    r"^duplicates-ignored\.csv$",
    r"^ignored-files\.csv$",
]

IGNORE_PATH_PARTS = {
    ".git", "__pycache__", ".trash", ".cache", "_manifest_reports",
}

MAX_CHUNK_CHARS = 1800
MIN_CHUNK_CHARS = 420
MAX_PREVIEW_SUMMARY = 360


@dataclass
class PageText:
    page: str
    text: str


@dataclass
class Candidate:
    path: Path
    rel_docs: Path
    rel_project: str
    sha256: str
    size: int


def strip_accents(text: str) -> str:
    return "".join(
        char for char in unicodedata.normalize("NFD", text)
        if unicodedata.category(char) != "Mn"
    )


def normalize_for_match(text: str) -> str:
    return "".join(
        char for char in unicodedata.normalize("NFD", str(text or "").lower())
        if unicodedata.category(char) != "Mn"
    )


SEARCH_STOP_WORDS = {
    "para", "com", "das", "dos", "uma", "por", "que", "seu", "sua",
    "de", "do", "da", "e", "o", "a", "as", "os", "em", "no", "na",
    "nos", "nas", "ao", "aos", "à", "às", "the", "and", "for",
}


def search_tokens(text: str) -> list[str]:
    normalized = re.sub(r"[^a-z0-9ç\s]", " ", normalize_for_match(text))
    seen: set[str] = set()
    tokens: list[str] = []
    for token in normalized.split():
        if len(token) <= 2 or token in SEARCH_STOP_WORDS or token in seen:
            continue
        seen.add(token)
        tokens.append(token)
    return tokens


def normalized_document_fields(title: str, summary: str, tags: list[str], kind: str, correspondent: str, fmt: str, status: str = "verified") -> dict:
    title_norm = normalize_for_match(title)
    summary_norm = normalize_for_match(summary)
    tags_norm = normalize_for_match(" ".join(tags))
    meta_norm = normalize_for_match(f"{kind} {correspondent} {fmt} {status}")
    return {
        "title": title_norm,
        "summary": summary_norm,
        "tags": tags_norm,
        "meta": meta_norm,
        "all": " ".join(part for part in (title_norm, summary_norm, tags_norm, meta_norm) if part),
        "tokens": search_tokens(f"{title} {summary} {' '.join(tags)} {kind} {correspondent} {fmt}"),
    }


def normalized_chunk_fields(heading: str, text: str, tags: list[str]) -> dict:
    heading_norm = normalize_for_match(heading)
    text_norm = normalize_for_match(text)
    tags_norm = normalize_for_match(" ".join(tags))
    return {
        "heading": heading_norm,
        "text": text_norm,
        "tags": tags_norm,
        "all": " ".join(part for part in (heading_norm, text_norm, tags_norm) if part),
        "tokens": search_tokens(f"{heading} {text} {' '.join(tags)}"),
    }


def normalize_spaces(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def compact(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def slugify(text: str) -> str:
    text = strip_accents(text).lower()
    text = re.sub(r"[^a-z0-9]+", "-", text).strip("-")
    return text or "documento"


def sha256_file(path: Path, chunk_size: int = 1024 * 1024) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            chunk = f.read(chunk_size)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def pretty_title(path: Path) -> str:
    stem = re.sub(r"[_-]+", " ", path.stem).strip()
    words: list[str] = []
    for word in stem.split():
        raw = strip_accents(word).lower()
        if raw in ACRONYMS:
            words.append(raw.upper())
        elif re.fullmatch(r"\d{4}(?:\.\d)?", word):
            words.append(word)
        elif re.fullmatch(r"\d+", word):
            words.append(word)
        else:
            words.append(word[:1].upper() + word[1:])
    return " ".join(words) or path.stem


def category_for(rel_docs: Path) -> str:
    if len(rel_docs.parts) <= 1:
        return "Documentos"
    first = rel_docs.parts[0]
    return CATEGORY_LABELS.get(first.lower(), re.sub(r"[-_]+", " ", first).title())


def infer_kind(text: str) -> str:
    hay = normalize_for_match(text)
    if "projeto pedagogico" in hay or re.search(r"\bppc\b", hay):
        return "PPC"
    if "matriz" in hay or "fluxograma" in hay:
        return "Matriz curricular"
    if "ementario" in hay or "bibliografia" in hay:
        return "Ementário"
    if "regulamento" in hay or "regimento" in hay:
        return "Regulamento"
    if "resolucao" in hay:
        return "Resolução"
    if "portaria" in hay:
        return "Portaria"
    if "lei" in hay or "ldb" in hay:
        return "Lei"
    if "diretriz" in hay:
        return "Diretriz"
    if "barema" in hay:
        return "Barema"
    if "formulario" in hay:
        return "Formulário"
    return "Documento"


def infer_correspondent(text: str) -> str:
    hay = normalize_for_match(text)
    if "consepe" in hay:
        return "CONSEPE"
    if "consup" in hay:
        return "CONSUP"
    if "concam" in hay:
        return "CONCAM"
    if "conaes" in hay:
        return "CONAES"
    if "cne" in hay or "mec" in hay or "ldb" in hay or "sinaes" in hay or "pne" in hay:
        return "MEC/CNE"
    if "nde" in hay or "nucleo docente estruturante" in hay:
        return "NDE"
    if "colegiado" in hay:
        return "Colegiado de Curso"
    if "ifba" in hay or "naes" in hay:
        return "IFBA"
    if "ppc" in hay or "matriz" in hay or "tcc" in hay or "estagio" in hay:
        return "Coordenação de Sistemas de Informação"
    return "A definir"


def infer_tags(title: str, category: str, kind: str, rel: str, extracted_text: str = "") -> list[str]:
    hay = normalize_for_match(f"{title} {category} {kind} {rel} {extracted_text[:4000]}")
    tags = [category, kind]
    rules = {
        "ppc": ["PPC", "projeto pedagógico", "currículo"],
        "projeto pedagogico": ["PPC", "projeto pedagógico", "currículo"],
        "matriz": ["matriz", "currículo", "disciplinas"],
        "fluxograma": ["fluxograma", "matriz", "currículo"],
        "optativa": ["optativas", "disciplinas"],
        "ementario": ["ementário", "bibliografia", "disciplinas"],
        "bibliografia": ["bibliografia", "disciplinas"],
        "migracao": ["migração", "equivalência", "currículo"],
        "equivalencia": ["equivalência", "aproveitamento"],
        "aproveitamento": ["aproveitamento", "disciplinas"],
        "tcc": ["TCC", "trabalho de conclusão"],
        "trabalho de conclusao": ["TCC", "trabalho de conclusão"],
        "estagio": ["estágio"],
        "extensao": ["extensão"],
        "curricularizacao": ["extensão", "curricularização"],
        "complementares": ["atividades complementares", "horas"],
        "libras": ["Libras", "acessibilidade"],
        "tea": ["TEA", "acessibilidade"],
        "deficiencia": ["acessibilidade", "inclusão"],
        "inclusao": ["inclusão", "acessibilidade"],
        "nde": ["NDE"],
        "nucleo docente estruturante": ["NDE"],
        "colegiado": ["colegiado"],
        "portaria": ["portaria"],
        "resolucao": ["resolução"],
        "regulamento": ["regulamento"],
        "regimento": ["regimento"],
        "naes": ["NAES", "normas acadêmicas"],
        "nome social": ["nome social", "inclusão"],
        "guarda religiosa": ["guarda religiosa", "atividades alternativas"],
    }
    for key, values in rules.items():
        if key in hay:
            tags.extend(values)
    tags.extend(re.findall(r"20\d{2}|19\d{2}", hay))

    result: list[str] = []
    seen: set[str] = set()
    for tag in tags:
        key = normalize_for_match(tag)
        if key and key not in seen:
            seen.add(key)
            result.append(tag)
    return result[:18]


def same_name_primary_exists(path: Path) -> bool:
    return any(path.with_suffix(ext).exists() for ext in PRIMARY_EXTENSIONS)


def is_ignored_name(path: Path) -> tuple[bool, str]:
    normalized = normalize_for_match(path.name)
    for pattern in IGNORE_NAME_PATTERNS:
        if re.fullmatch(pattern, normalized):
            return True, f"ignored pattern: {pattern}"
    return False, ""


def should_ignore(path: Path) -> tuple[bool, str]:
    if not path.is_file():
        return True, "not a file"

    rel_parts = [normalize_for_match(part) for part in path.relative_to(DOCS_DIR).parts]
    if any(part in IGNORE_PATH_PARTS for part in rel_parts):
        return True, "ignored folder"

    ignored_name, reason = is_ignored_name(path)
    if ignored_name:
        return True, reason

    if path.suffix.lower() not in SUPPORTED:
        return True, "unsupported extension"

    # Ignore TXT/MD sidecar files when a same-name primary document exists.
    # Their text can still be used for that document.
    if path.suffix.lower() in {".txt", ".md"} and same_name_primary_exists(path):
        return True, "sidecar for same-name primary document"

    return False, ""


def read_sidecar(path: Path) -> list[PageText]:
    candidates = [path.with_suffix(".txt"), path.with_suffix(".md")]
    for candidate in candidates:
        if candidate.exists() and candidate != path:
            ignored_name, _ = is_ignored_name(candidate)
            if ignored_name:
                continue
            try:
                text = normalize_spaces(candidate.read_text(encoding="utf-8", errors="ignore"))
                if text:
                    return [PageText("—", text)]
            except Exception:
                pass
    return []


def extract_pdf_pymupdf(path: Path) -> list[PageText]:
    try:
        import fitz  # PyMuPDF
    except Exception:
        return []
    pages: list[PageText] = []
    try:
        with fitz.open(path) as doc:
            for i, page in enumerate(doc, start=1):
                text = normalize_spaces(page.get_text("text") or "")
                if text:
                    pages.append(PageText(str(i), text))
    except Exception as exc:
        print(f"[aviso] PyMuPDF falhou em {path.name}: {exc}", file=sys.stderr)
    return pages


def extract_pdf_pypdf(path: Path) -> list[PageText]:
    try:
        from pypdf import PdfReader
    except Exception:
        return []
    pages: list[PageText] = []
    try:
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, start=1):
            text = normalize_spaces(page.extract_text() or "")
            if text:
                pages.append(PageText(str(i), text))
    except Exception as exc:
        print(f"[aviso] pypdf falhou em {path.name}: {exc}", file=sys.stderr)
    return pages


def extract_pdf_pdftotext(path: Path) -> list[PageText]:
    if not shutil.which("pdftotext"):
        return []
    try:
        output = subprocess.check_output(
            ["pdftotext", "-layout", str(path), "-"],
            stderr=subprocess.DEVNULL,
            timeout=60,
        )
        text = output.decode("utf-8", errors="ignore")
        parts = [normalize_spaces(part) for part in text.split("\f")]
        return [PageText(str(i), part) for i, part in enumerate(parts, start=1) if part]
    except Exception as exc:
        print(f"[aviso] pdftotext falhou em {path.name}: {exc}", file=sys.stderr)
        return []


def extract_docx(path: Path) -> list[PageText]:
    try:
        import docx
    except Exception:
        return []
    try:
        document = docx.Document(str(path))
        text = "\n".join(p.text for p in document.paragraphs if p.text)
        return [PageText("—", normalize_spaces(text))] if text.strip() else []
    except Exception as exc:
        print(f"[aviso] docx falhou em {path.name}: {exc}", file=sys.stderr)
        return []


def extract_xlsx(path: Path) -> list[PageText]:
    try:
        import openpyxl
    except Exception:
        return []
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        pages: list[PageText] = []
        for sheet in wb.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    rows.append(" | ".join(values))
            text = normalize_spaces("\n".join(rows))
            if text:
                pages.append(PageText(sheet.title, text))
        return pages
    except Exception as exc:
        print(f"[aviso] xlsx falhou em {path.name}: {exc}", file=sys.stderr)
        return []


def extract_pptx(path: Path) -> list[PageText]:
    try:
        from pptx import Presentation
    except Exception:
        return []
    try:
        prs = Presentation(str(path))
        pages: list[PageText] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            text = normalize_spaces("\n".join(parts))
            if text:
                pages.append(PageText(str(i), text))
        return pages
    except Exception as exc:
        print(f"[aviso] pptx falhou em {path.name}: {exc}", file=sys.stderr)
        return []


def extract_plain(path: Path) -> list[PageText]:
    try:
        text = normalize_spaces(path.read_text(encoding="utf-8", errors="ignore"))
        return [PageText("—", text)] if text else []
    except Exception:
        return []


def parse_pdf_date(value: object) -> str:
    """Return YYYY-MM-DD from PDF/XMP/Office-style dates, or empty string."""
    if not value:
        return ""
    raw = str(value).strip()
    # Common PDF metadata: D:YYYYMMDDHHmmSS-03'00'
    match = re.search(r"D?:?\s*(\d{4})(\d{2})(\d{2})", raw)
    if match:
        return f"{match.group(1)}-{match.group(2)}-{match.group(3)}"
    # ISO-like: 2026-01-23 or 2026-01-23T10:00:00Z
    match = re.search(r"(\d{4})-(\d{2})-(\d{2})", raw)
    if match:
        return f"{match.group(1)}-{match.group(2)}-{match.group(3)}"
    # Brazilian date: 23/01/2026
    match = re.search(r"\b(\d{1,2})[/-](\d{1,2})[/-](20\d{2}|19\d{2})\b", raw)
    if match:
        day, month, year = match.groups()
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"
    return ""


def date_from_filename(path: Path) -> str:
    text = normalize_for_match(path.stem)
    # yyyy-mm-dd / yyyy_mm_dd / yyyy mm dd
    match = re.search(r"\b(20\d{2}|19\d{2})[-_ .](\d{1,2})[-_ .](\d{1,2})\b", text)
    if match:
        year, month, day = match.groups()
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"
    # dd-mm-yyyy / dd_mm_yyyy
    match = re.search(r"\b(\d{1,2})[-_ .](\d{1,2})[-_ .](20\d{2}|19\d{2})\b", text)
    if match:
        day, month, year = match.groups()
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"
    return ""


def date_from_text(text: str) -> str:
    sample = text[:12000]
    # Prefer explicit document dates near common words.
    patterns = [
        r"(?:de|em|data|publicad[ao]|aprovad[ao]|portaria|resolu[cç][aã]o)[^\n]{0,80}?\b(\d{1,2})\s+de\s+([a-zç]+)\s+de\s+(20\d{2}|19\d{2})\b",
        r"\b(\d{1,2})\s+de\s+([a-zç]+)\s+de\s+(20\d{2}|19\d{2})\b",
        r"\b(\d{1,2})[/-](\d{1,2})[/-](20\d{2}|19\d{2})\b",
    ]
    months = {
        "janeiro": 1, "fevereiro": 2, "marco": 3, "março": 3, "abril": 4,
        "maio": 5, "junho": 6, "julho": 7, "agosto": 8, "setembro": 9,
        "outubro": 10, "novembro": 11, "dezembro": 12,
    }
    normalized = normalize_for_match(sample)
    for pattern in patterns[:2]:
        match = re.search(pattern, normalized, flags=re.I)
        if match:
            day, month_name, year = match.groups()[-3:]
            month = months.get(month_name)
            if month:
                return f"{int(year):04d}-{month:02d}-{int(day):02d}"
    match = re.search(patterns[2], sample)
    if match:
        day, month, year = match.groups()
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"
    return ""


def office_created_date(path: Path) -> str:
    """Extract created/modified date from OOXML files without changing file timestamps."""
    try:
        with zipfile.ZipFile(path) as zf:
            if "docProps/core.xml" not in zf.namelist():
                return ""
            root = ET.fromstring(zf.read("docProps/core.xml"))
            ns = {
                "dcterms": "http://purl.org/dc/terms/",
                "dc": "http://purl.org/dc/elements/1.1/",
            }
            for tag in ["dcterms:created", "dcterms:modified", "dc:date"]:
                node = root.find(tag, ns)
                if node is not None and node.text:
                    parsed = parse_pdf_date(node.text)
                    if parsed:
                        return parsed
    except Exception:
        return ""
    return ""


def xmp_created_date_from_pdf(path: Path) -> str:
    try:
        import fitz
        with fitz.open(path) as doc:
            xmp = doc.get_xml_metadata() or ""
        if not xmp:
            return ""
        for tag in ["xmp:CreateDate", "xmp:ModifyDate", "pdf:ModDate", "dc:date"]:
            match = re.search(rf"<{re.escape(tag)}[^>]*>(.*?)</{re.escape(tag)}>", xmp, flags=re.I | re.S)
            if match:
                parsed = parse_pdf_date(match.group(1))
                if parsed:
                    return parsed
        parsed = parse_pdf_date(xmp)
        return parsed
    except Exception:
        return ""


def document_metadata(path: Path, extracted_pages: list[PageText] | None = None) -> dict:
    suffix = path.suffix.lower()
    full_text = " ".join(page.text for page in (extracted_pages or []))
    meta = {
        "pageCount": len(extracted_pages or []) or "",
        # Do not use the filesystem modified time here: it becomes the day you copied/downloaded
        # the file, which made every document appear as the current date on the website.
        "createdDate": "",
    }

    if suffix == ".pdf":
        try:
            import fitz  # PyMuPDF
            with fitz.open(path) as doc:
                meta["pageCount"] = len(doc)
                metadata = doc.metadata or {}
                created = metadata.get("creationDate") or metadata.get("modDate")
                meta["createdDate"] = parse_pdf_date(created) or xmp_created_date_from_pdf(path) or date_from_filename(path) or date_from_text(full_text)
                return meta
        except Exception:
            pass
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            meta["pageCount"] = len(reader.pages)
            metadata = reader.metadata or {}
            created = metadata.get("/CreationDate") or metadata.get("/ModDate")
            meta["createdDate"] = parse_pdf_date(created) or date_from_filename(path) or date_from_text(full_text)
            return meta
        except Exception:
            meta["createdDate"] = date_from_filename(path) or date_from_text(full_text)
            return meta

    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        meta["pageCount"] = 1
        meta["createdDate"] = date_from_filename(path)
    elif suffix in {".docx", ".xlsx", ".pptx"}:
        meta["createdDate"] = office_created_date(path) or date_from_filename(path) or date_from_text(full_text)
        if suffix == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                meta["pageCount"] = len(wb.worksheets)
            except Exception:
                pass
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                meta["pageCount"] = len(Presentation(str(path)).slides)
            except Exception:
                pass
        elif suffix == ".docx" and not meta["pageCount"]:
            meta["pageCount"] = 1
    elif suffix in {".txt", ".md"}:
        meta["pageCount"] = 1
        meta["createdDate"] = date_from_filename(path) or date_from_text(full_text)

    return meta

def extract_pages(path: Path) -> tuple[list[PageText], str]:
    suffix = path.suffix.lower()

    sidecar = read_sidecar(path)
    if sidecar:
        return sidecar, "sidecar"

    if suffix == ".pdf":
        pages = extract_pdf_pymupdf(path)
        method = "pymupdf"
        if not pages:
            pages = extract_pdf_pypdf(path)
            method = "pypdf"
        if not pages:
            pages = extract_pdf_pdftotext(path)
            method = "pdftotext"
        return pages, method if pages else "no_text"

    if suffix in {".txt", ".md"}:
        return extract_plain(path), "plain"
    if suffix == ".docx":
        return extract_docx(path), "docx"
    if suffix == ".xlsx":
        return extract_xlsx(path), "xlsx"
    if suffix == ".pptx":
        return extract_pptx(path), "pptx"

    return [], "unsupported_text_extraction"


def split_text_into_chunks(text: str, max_chars: int = MAX_CHUNK_CHARS) -> list[str]:
    text = normalize_spaces(text)
    if not text:
        return []
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    if len(paragraphs) <= 1:
        paragraphs = [p.strip() for p in re.split(r"(?<=[.!?;:])\s+(?=[A-ZÁÉÍÓÚÃÕÂÊÔÇ0-9])", text) if p.strip()]

    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        paragraph = compact(paragraph)
        if not paragraph:
            continue
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current.strip())
                current = ""
            for start in range(0, len(paragraph), max_chars):
                chunks.append(paragraph[start:start + max_chars].strip())
            continue
        if current and len(current) + len(paragraph) + 1 > max_chars and len(current) >= MIN_CHUNK_CHARS:
            chunks.append(current.strip())
            current = paragraph
        else:
            current = f"{current} {paragraph}".strip()
    if current:
        chunks.append(current.strip())
    return chunks


def make_chunks(doc_id: str, title: str, pages: list[PageText], tags: list[str]) -> list[dict]:
    chunks: list[dict] = []
    for page in pages:
        for idx, chunk_text in enumerate(split_text_into_chunks(page.text), start=1):
            heading = title if idx == 1 else f"{title} — trecho {idx}"
            chunks.append({
                "id": f"{doc_id}-p{slugify(str(page.page))}-{idx}",
                "page": page.page,
                "heading": heading,
                "semanticTags": tags,
                "text": chunk_text,
                "normalized": normalized_chunk_fields(heading, chunk_text, tags),
            })
    return chunks


def first_summary(text: str, fallback: str) -> str:
    clean = compact(text)
    if not clean:
        return fallback
    return clean[:MAX_PREVIEW_SUMMARY] + ("…" if len(clean) > MAX_PREVIEW_SUMMARY else "")


def preferred_candidate(candidates: list[Candidate]) -> Candidate:
    """Choose which duplicate file should represent one hash in the manifest."""
    def score(c: Candidate) -> tuple[int, int, int, str]:
        parts = [part.lower() for part in c.rel_docs.parts]
        first = parts[0] if parts else ""
        in_known_category = 1 if first in KNOWN_CATEGORY_FOLDERS else 0
        in_subfolder = 1 if len(parts) > 1 else 0
        bad_copy_name = 1 if re.search(r"\b(copy|copia|cópia|duplicate|duplicado)\b", normalize_for_match(c.path.name)) else 0
        # Higher is better, except bad_copy_name and path string.
        return (in_known_category, in_subfolder, -bad_copy_name, c.rel_project)

    return sorted(candidates, key=score, reverse=True)[0]


def collect_candidates() -> tuple[list[Candidate], list[dict]]:
    candidates: list[Candidate] = []
    ignored: list[dict] = []

    if not DOCS_DIR.exists():
        DOCS_DIR.mkdir(parents=True, exist_ok=True)

    for path in sorted(DOCS_DIR.rglob("*")):
        ignore, reason = should_ignore(path)
        if ignore:
            if path.is_file() and reason not in {"unsupported extension", "not a file"}:
                ignored.append({
                    "path": path.relative_to(ROOT).as_posix(),
                    "reason": reason,
                })
            continue
        try:
            candidates.append(Candidate(
                path=path,
                rel_docs=path.relative_to(DOCS_DIR),
                rel_project=path.relative_to(ROOT).as_posix(),
                sha256=sha256_file(path),
                size=path.stat().st_size,
            ))
        except Exception as exc:
            ignored.append({
                "path": path.relative_to(ROOT).as_posix(),
                "reason": f"could not hash file: {exc}",
            })

    return candidates, ignored


def deduplicate_candidates(candidates: list[Candidate]) -> tuple[list[Candidate], list[dict]]:
    by_hash: dict[str, list[Candidate]] = {}
    for candidate in candidates:
        by_hash.setdefault(candidate.sha256, []).append(candidate)

    unique: list[Candidate] = []
    duplicates: list[dict] = []

    for sha, group in by_hash.items():
        if len(group) == 1:
            unique.append(group[0])
            continue
        kept = preferred_candidate(group)
        unique.append(kept)
        for duplicate in group:
            if duplicate.path == kept.path:
                continue
            duplicates.append({
                "duplicatePath": duplicate.rel_project,
                "keptPath": kept.rel_project,
                "sha256": sha,
                "size": duplicate.size,
            })

    unique.sort(key=lambda c: c.rel_project)
    duplicates.sort(key=lambda row: row["duplicatePath"])
    return unique, duplicates


def unique_doc_id(base: str, used_ids: set[str]) -> str:
    if base not in used_ids:
        used_ids.add(base)
        return base
    i = 2
    while f"{base}-{i}" in used_ids:
        i += 1
    final = f"{base}-{i}"
    used_ids.add(final)
    return final


def build_document(candidate: Candidate, used_ids: set[str]) -> dict:
    path = candidate.path
    rel = candidate.rel_project
    rel_docs = candidate.rel_docs
    category = category_for(rel_docs)
    title = pretty_title(path)
    pages, extraction_method = extract_pages(path)
    file_meta = document_metadata(path, pages)
    full_text = compact(" ".join(page.text for page in pages))
    kind = infer_kind(f"{title} {category} {rel} {full_text[:4000]}")
    correspondent = infer_correspondent(f"{title} {category} {rel} {full_text[:5000]}")
    fmt = SUPPORTED.get(path.suffix.lower(), path.suffix.upper().replace(".", ""))
    tags = infer_tags(title, category, kind, rel, full_text)
    base_id = f"doc-{slugify(str(rel_docs.with_suffix('')))}"
    doc_id = unique_doc_id(base_id, used_ids)

    if pages:
        chunks = make_chunks(doc_id, title, pages, tags)
    else:
        tags = list(dict.fromkeys([*tags, "sem texto extraível"]))
        fallback_text = f"{title}. {kind} em {category}. Arquivo disponível em {rel}. Atenção: não foi possível extrair texto automaticamente deste arquivo; se for PDF escaneado, gere OCR ou adicione um .txt com o mesmo nome."
        chunks = [{
            "id": f"{doc_id}-arquivo",
            "page": "—",
            "heading": title,
            "semanticTags": tags,
            "text": fallback_text,
            "normalized": normalized_chunk_fields(title, fallback_text, tags),
        }]

    summary = first_summary(full_text, f"{kind} em {category}.")
    thumbnail_meta = generate_pdf_thumbnails(path, doc_id)
    normalized_fields = normalized_document_fields(title, summary, tags, kind, correspondent, fmt)
    return {
        "id": doc_id,
        "title": title,
        "category": category,
        "group": category,
        "kind": kind,
        "documentType": kind,
        "correspondent": correspondent,
        "fileFormat": fmt,
        "pdfUrl": rel,
        "sourceUrl": rel,
        "sourceLabel": correspondent or "IFBA / fonte institucional",
        "reviewedDate": "",
        "validityStatus": "A conferir",
        "supersededBy": "",
        "tags": tags,
        "summary": summary,
        "normalized": normalized_fields,
        "pageCount": file_meta.get("pageCount") or "",
        "createdDate": file_meta.get("createdDate") or "",
        "contentLength": len(full_text),
        "indexed": bool(full_text),
        "extractionMethod": extraction_method,
        "sha256": candidate.sha256,
        "size": candidate.size,
        **thumbnail_meta,
        "chunks": chunks,
    }



def generate_pdf_thumbnails(path: Path, doc_id: str) -> dict:
    """Generate responsive first-page WebP thumbnails when optional libraries exist."""
    if path.suffix.lower() != ".pdf":
        return {}
    try:
        import fitz
        from PIL import Image
        THUMBNAIL_DIR.mkdir(parents=True, exist_ok=True)
        outputs = []
        first_dimensions = None
        with fitz.open(path) as pdf:
            if len(pdf) < 1:
                return {}
            page = pdf[0]
            rect = page.rect
            for width in THUMBNAIL_SIZES:
                scale = width / max(1.0, rect.width)
                pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                target = THUMBNAIL_DIR / f"{doc_id}-{width}.webp"
                image.save(target, "WEBP", quality=76, method=6)
                outputs.append((target.relative_to(ROOT).as_posix(), width))
                if width == 320:
                    first_dimensions = (pix.width, pix.height)
        if not outputs:
            return {}
        preferred = next((path for path, width in outputs if width == 320), outputs[0][0])
        width, height = first_dimensions or (320, 452)
        return {
            "thumbnailUrl": preferred,
            "thumbnailSrcset": ", ".join(f"{path} {size}w" for path, size in outputs),
            "thumbnailWidth": width,
            "thumbnailHeight": height,
        }
    except Exception as exc:
        print(f"Thumbnail skipped for {path.relative_to(ROOT)}: {exc}", file=sys.stderr)
        return {}

def write_csv(path: Path, rows: list[dict], fieldnames: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def main() -> None:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    THUMBNAIL_DIR.mkdir(parents=True, exist_ok=True)
    # Rebuild outputs from the current manifest source so removed/renamed PDFs
    # cannot leave orphaned thumbnails in future releases.
    for stale in THUMBNAIL_DIR.glob("*.webp"):
        stale.unlink(missing_ok=True)

    candidates, ignored = collect_candidates()
    unique_candidates, duplicates = deduplicate_candidates(candidates)

    used_ids: set[str] = set()
    docs = [build_document(candidate, used_ids) for candidate in unique_candidates]
    indexed = sum(1 for doc in docs if doc.get("indexed"))

    generated_at = datetime.now(timezone.utc).isoformat()
    manifest = {
        "generatedAt": generated_at,
        "count": len(docs),
        "indexedCount": indexed,
        "unindexedCount": len(docs) - indexed,
        "duplicateFilesIgnored": len(duplicates),
        "ignoredHelperFiles": len(ignored),
        "summaryUrl": "documents/manifest-summary.json",
        "searchIndexUrl": "documents/search-index.json",
        "documents": docs,
    }
    MANIFEST_JSON.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    summary_documents = []
    search_documents = []
    vocabulary: dict[str, dict] = {}

    def add_vocabulary(display: str, weight: int) -> None:
        for word in re.findall(r"[A-Za-zÀ-ÖØ-öø-ÿÇç0-9]+", str(display or "")):
            key = normalize_for_match(word)
            if len(key) < 4 or key in SEARCH_STOP_WORDS or key.isdigit():
                continue
            current = vocabulary.setdefault(key, {"key": key, "display": word, "weight": 0})
            current["weight"] += weight
            if word != normalize_for_match(word) and current["display"] == normalize_for_match(current["display"]):
                current["display"] = word

    for doc in docs:
        summary_doc = {key: value for key, value in doc.items() if key != "chunks"}
        summary_doc["hasPassages"] = bool(doc.get("chunks"))
        summary_doc["passageCount"] = len(doc.get("chunks", []))
        summary_documents.append(summary_doc)
        search_documents.append({
            "id": doc["id"],
            "normalized": doc.get("normalized", {}),
            "chunks": doc.get("chunks", []),
        })
        add_vocabulary(doc.get("title", ""), 12)
        add_vocabulary(" ".join(doc.get("tags", [])), 7)
        add_vocabulary(f"{doc.get('kind', '')} {doc.get('correspondent', '')}", 5)
        add_vocabulary(doc.get("summary", ""), 2)

    summary_manifest = {
        "generatedAt": generated_at,
        "count": len(docs),
        "indexedCount": indexed,
        "unindexedCount": len(docs) - indexed,
        "searchIndexUrl": "documents/search-index.json",
        "documents": summary_documents,
    }
    search_index = {
        "generatedAt": generated_at,
        "version": 1,
        "documentCount": len(search_documents),
        "passageCount": sum(len(item["chunks"]) for item in search_documents),
        "vocabulary": sorted(vocabulary.values(), key=lambda item: (-item["weight"], item["key"])),
        "documents": search_documents,
    }
    MANIFEST_SUMMARY_JSON.write_text(json.dumps(summary_manifest, ensure_ascii=False, separators=(",", ":")), encoding="utf-8")
    SEARCH_INDEX_JSON.write_text(json.dumps(search_index, ensure_ascii=False, separators=(",", ":")), encoding="utf-8")

    write_csv(MANIFEST_CSV, [
        {
            "id": doc["id"],
            "title": doc["title"],
            "category": doc["category"],
            "kind": doc["kind"],
            "correspondent": doc["correspondent"],
            "fileFormat": doc["fileFormat"],
            "path": doc["pdfUrl"],
            "thumbnailUrl": doc.get("thumbnailUrl", ""),
            "thumbnailSrcset": doc.get("thumbnailSrcset", ""),
            "pageCount": doc.get("pageCount", ""),
            "createdDate": doc.get("createdDate", ""),
            "indexed": doc.get("indexed", False),
            "contentLength": doc.get("contentLength", 0),
            "extractionMethod": doc.get("extractionMethod", ""),
            "chunks": len(doc.get("chunks", [])),
            "sha256": doc.get("sha256", ""),
            "tags": "; ".join(doc["tags"]),
            "summary": doc["summary"],
        }
        for doc in docs
    ], [
        "id", "title", "category", "kind", "correspondent", "fileFormat", "path",
        "thumbnailUrl", "thumbnailSrcset", "pageCount", "createdDate", "indexed", "contentLength", "extractionMethod", "chunks", "sha256", "tags", "summary",
    ])

    write_csv(DUPLICATES_CSV, duplicates, ["duplicatePath", "keptPath", "sha256", "size"])
    write_csv(IGNORED_CSV, ignored, ["path", "reason"])

    print(f"Generated {MANIFEST_JSON.relative_to(ROOT)} with {len(docs)} document(s).")
    print(f"Generated lightweight metadata: {MANIFEST_SUMMARY_JSON.relative_to(ROOT)}")
    print(f"Generated deferred passage index: {SEARCH_INDEX_JSON.relative_to(ROOT)}")
    print(f"Indexed content in {indexed}/{len(docs)} document(s).")
    print(f"Ignored helper files: {len(ignored)}")
    print(f"Duplicate files ignored: {len(duplicates)}")
    if duplicates:
        print(f"Duplicate report: {DUPLICATES_CSV.relative_to(ROOT)}")
    if ignored:
        print(f"Ignored-file report: {IGNORED_CSV.relative_to(ROOT)}")
    if indexed < len(docs):
        print("Some files have no extractable text. If they are scanned PDFs, add OCR or a same-name .txt sidecar.")
    print(f"Generated {MANIFEST_CSV.relative_to(ROOT)} for review.")


if __name__ == "__main__":
    main()
